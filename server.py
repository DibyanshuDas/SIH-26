"""
SKILL-STAT Web Server (Port 8001)
Serves the AI-Enabled Official Statistical Learning, Competency Intelligence & Assessment Platform.
Provides full REST APIs for:
- Learner Digital Competency Passport & Skill Gap Profiling
- Dual-Track iGOT Karmayogi & NSSTA TPAC Recommendation Engine
- AI Generative Assessment Engine (MCQ/Quiz generation from text/documents)
- Instant Quiz Evaluation, Explanations & Live Competency Score Uplift
- MoSPI Administrative Cadre Analytics, Division Heatmaps & Predictive Capacity Planning
- RAG Statistical Learning Assistant
"""

import http.server
import socketserver
import json
import os
import urllib.parse
import time
import random
import tempfile
import io
from datetime import datetime

PORT = int(os.environ.get("PORT", 8050))
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DASHBOARD_DIR = os.path.join(BASE_DIR, "dashboard")
DATA_DIR = os.path.join(BASE_DIR, "data")
DASHBOARD_DATA_DIR = os.path.join(DASHBOARD_DIR, "data")

from competency_engine import CompetencyEngine, COMPETENCY_FRAMEWORK, CADRE_ROLE_REQUIREMENTS, DIVISIONS
from recommendation_engine import RecommendationEngine
from assessment_generator import AssessmentEngine

# Initialize In-Memory Engines
competency_engine = CompetencyEngine()
recommendation_engine = RecommendationEngine()
assessment_engine = AssessmentEngine()

# Initialize data if not already existing
if not os.path.exists(os.path.join(DASHBOARD_DATA_DIR, "primary_learner.json")):
    competency_engine.generate_officer_profiles(2850)
    competency_engine.save_all_data()
    recommendation_engine.save_all_catalogs()
    assessment_engine.save_all_data()

class KashyapRequestHandler(http.server.SimpleHTTPRequestHandler):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, directory=DASHBOARD_DIR, **kwargs)

    def do_GET(self):
        parsed_url = urllib.parse.urlparse(self.path)
        path = parsed_url.path
        query_params = urllib.parse.parse_qs(parsed_url.query)

        # 1. API: Competency Framework Taxonomy
        if path == "/api/framework":
            self.send_json_response(COMPETENCY_FRAMEWORK)
            return

        # 2. API: Learner Profile
        elif path == "/api/learner-profile":
            officer_id = query_params.get("id", [None])[0]
            role_filter = query_params.get("role", [None])[0]
            
            profile = self.get_officer_profile(officer_id, role_filter)
            self.send_json_response(profile)
            return

        # 3. API: Search / List Officers
        elif path == "/api/officers":
            division = query_params.get("division", [None])[0]
            cadre = query_params.get("cadre", [None])[0]
            search = query_params.get("q", [None])[0]
            limit = int(query_params.get("limit", [50])[0])
            self.send_json_response(self.search_officers(division, cadre, search, limit))
            return

        # 4. API: Skill Gaps
        elif path == "/api/skill-gaps":
            officer_id = query_params.get("id", [None])[0]
            profile = self.get_officer_profile(officer_id)
            self.send_json_response({
                "officer_id": profile["officer_id"],
                "overall_competency_index": profile["overall_competency_index"],
                "domain_scores": profile.get("domain_scores", {}),
                "skill_gaps": profile.get("skill_gaps", {}),
                "top_priority_gaps": profile.get("top_priority_gaps", [])
            })
            return

        # 5. API: Recommendations (iGOT + NSSTA)
        elif path == "/api/recommendations":
            officer_id = query_params.get("id", [None])[0]
            profile = self.get_officer_profile(officer_id)
            recs = recommendation_engine.recommend_for_officer(profile)
            self.send_json_response(recs)
            return

        # 6. API: iGOT Course Catalog
        elif path == "/api/igot/courses":
            category = query_params.get("category", [None])[0]
            search = query_params.get("q", [None])[0]
            courses = recommendation_engine.igot_courses
            if category and category != "All":
                courses = [c for c in courses if category in c.get("primary_competency", "") or any(category in t for t in c.get("competency_tags", []))]
            if search:
                s_lower = search.lower()
                courses = [c for c in courses if s_lower in c["title"].lower() or s_lower in c["description"].lower()]
            self.send_json_response(courses)
            return

        # 7. API: NSSTA TPAC Programmes
        elif path == "/api/nssta/programmes":
            self.send_json_response(recommendation_engine.tpac_programmes)
            return

        # 8. API: Learning Materials
        elif path == "/api/materials":
            self.send_json_response(assessment_engine.get_all_materials())
            return

        # 9. API: Assessment Bank
        elif path == "/api/assessments":
            self.send_json_response(assessment_engine.get_all_assessments())
            return

        # 10. API: Single Assessment Detail
        elif path == "/api/assessments/detail":
            asm_id = query_params.get("id", [None])[0]
            asm = assessment_engine.get_assessment_by_id(asm_id)
            if asm:
                self.send_json_response(asm)
            else:
                self.send_json_response({"error": "Assessment not found"})
            return

        # 11. API: Administrative Analytics
        elif path == "/api/admin/analytics":
            analytics_file = os.path.join(DASHBOARD_DATA_DIR, "administrative_analytics.json")
            if os.path.exists(analytics_file):
                with open(analytics_file, "r") as f:
                    self.send_json_response(json.load(f))
            else:
                self.send_json_response(competency_engine.compute_administrative_analytics())
            return

        # Default static file handler
        return super().do_GET()

    def do_POST(self):
        parsed_url = urllib.parse.urlparse(self.path)
        path = parsed_url.path

        content_length = int(self.headers.get("Content-Length", 0))
        content_type = self.headers.get("Content-Type", "")

        # Handle file upload (multipart/form-data)
        if path == "/api/assessments/upload-material":
            self.handle_file_upload(content_length, content_type)
            return

        body_bytes = self.rfile.read(content_length) if content_length > 0 else b"{}"
        try:
            payload = json.loads(body_bytes.decode("utf-8"))
        except Exception:
            payload = {}

        # 1. API: Enrol in iGOT Course
        if path == "/api/igot/enrol":
            course_id = payload.get("course_id")
            officer_id = payload.get("officer_id", "OFF-ISS-2026-HQ")
            res = self.handle_course_enrolment(officer_id, course_id)
            self.send_json_response(res)
            return

        # 2. API: Generate Assessment from Uploaded Text/Document
        elif path == "/api/assessments/generate":
            title = payload.get("title", "Uploaded Statistical Guidelines")
            content = payload.get("content", "")
            target_comp = payload.get("target_competency", "STAT-01")
            num_q = int(payload.get("num_questions", 5))

            if not content.strip():
                # Pick from default material if empty
                mats = list(assessment_engine.materials.values())
                content = mats[0]["content"]
                title = mats[0]["title"]

            generated_asm = assessment_engine.generate_assessment_from_text(title, content, target_comp, num_q)
            self.send_json_response({"success": True, "assessment": generated_asm})
            return

        # 3. API: NSSTA Nominate
        elif path == "/api/nssta/nominate":
            program_id = payload.get("program_id")
            officer_id = payload.get("officer_id", "OFF-ISS-2026-HQ")
            res = self.handle_nssta_nomination(officer_id, program_id)
            self.send_json_response(res)
            return

        # 3. API: Submit Assessment Answers
        elif path == "/api/assessments/submit":
            asm_id = payload.get("assessment_id")
            answers = payload.get("answers", {})  # { "Q-SNA-01": 1, ... }
            time_spent = payload.get("time_spent_seconds", 180)
            officer_id = payload.get("officer_id", "OFF-ISS-2026-HQ")

            eval_result = assessment_engine.evaluate_submission(asm_id, answers, time_spent)
            
            # If passed, update officer's competency and karma
            if eval_result.get("passed"):
                self.apply_assessment_uplift(officer_id, eval_result)

            self.send_json_response(eval_result)
            return

        # 4. API: AI Karmayogi Statistical Assistant RAG Query
        elif path == "/api/assistant/query":
            user_prompt = payload.get("query", "")
            officer_id = payload.get("officer_id", "OFF-ISS-2026-HQ")
            response = self.handle_assistant_query(user_prompt, officer_id)
            self.send_json_response(response)
            return

        self.send_response(404)
        self.end_headers()

    def send_json_response(self, data):
        self.send_response(200)
        self.send_header("Content-Type", "application/json")
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "GET, POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
        self.send_header("Cache-Control", "no-cache")
        self.end_headers()
        self.wfile.write(json.dumps(data).encode("utf-8"))

    def handle_file_upload(self, content_length, content_type):
        """Handle multipart file upload and extract text from PDF, DOCX, PPTX, TXT."""
        try:
            body_bytes = self.rfile.read(content_length) if content_length > 0 else b""

            # Parse multipart boundary
            if "boundary=" not in content_type:
                self.send_json_response({"success": False, "error": "Invalid upload format"})
                return

            boundary = content_type.split("boundary=")[1].strip()
            if boundary.startswith('"') and boundary.endswith('"'):
                boundary = boundary[1:-1]

            # Parse multipart parts
            parts = body_bytes.split(("--" + boundary).encode())
            file_data = None
            filename = ""

            for part in parts:
                if b"filename=" in part:
                    # Extract filename
                    header_section = part.split(b"\r\n\r\n", 1)
                    if len(header_section) < 2:
                        continue
                    headers_str = header_section[0].decode("utf-8", errors="ignore")
                    file_content = header_section[1]
                    # Remove trailing boundary markers
                    if file_content.endswith(b"\r\n"):
                        file_content = file_content[:-2]

                    # Extract filename from Content-Disposition header
                    for line in headers_str.split("\r\n"):
                        if "filename=" in line:
                            fname_part = line.split('filename="')[1] if 'filename="' in line else ""
                            filename = fname_part.split('"')[0] if fname_part else "uploaded_file"
                            break

                    file_data = file_content
                    break

            if not file_data or not filename:
                self.send_json_response({"success": False, "error": "No file found in upload"})
                return

            ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
            extracted_text = ""
            page_info = ""

            # Save to temp file for processing
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(file_data)
                tmp_path = tmp.name

            try:
                if ext == ".txt":
                    extracted_text = file_data.decode("utf-8", errors="ignore")

                elif ext == ".pdf":
                    try:
                        import PyPDF2
                        reader = PyPDF2.PdfReader(tmp_path)
                        pages = []
                        for page in reader.pages:
                            text = page.extract_text()
                            if text:
                                pages.append(text)
                        extracted_text = "\n\n".join(pages)
                        page_info = f"({len(reader.pages)} pages extracted)"
                    except ImportError:
                        extracted_text = "[PDF extraction requires PyPDF2. Install: pip install PyPDF2]"

                elif ext in [".docx", ".doc"]:
                    try:
                        import docx
                        doc = docx.Document(tmp_path)
                        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                        extracted_text = "\n\n".join(paragraphs)
                        page_info = f"({len(paragraphs)} paragraphs extracted)"
                    except ImportError:
                        extracted_text = "[DOCX extraction requires python-docx. Install: pip install python-docx]"

                elif ext == ".pptx":
                    try:
                        from pptx import Presentation
                        prs = Presentation(tmp_path)
                        slides_text = []
                        for slide_num, slide in enumerate(prs.slides, 1):
                            slide_content = []
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text.strip():
                                    slide_content.append(shape.text)
                            if slide_content:
                                slides_text.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_content))
                        extracted_text = "\n\n".join(slides_text)
                        page_info = f"({len(prs.slides)} slides extracted)"
                    except ImportError:
                        extracted_text = "[PPTX extraction requires python-pptx. Install: pip install python-pptx]"
                else:
                    extracted_text = "[Unsupported file format]"

            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass

            if not extracted_text.strip():
                extracted_text = "[No readable text could be extracted from this file.]"

            self.send_json_response({
                "success": True,
                "extracted_text": extracted_text[:50000],  # Cap at 50k chars
                "filename": filename,
                "pages": page_info
            })

        except Exception as e:
            print(f"File upload error: {e}")
            self.send_json_response({"success": False, "error": str(e)})

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "GET, POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
        self.end_headers()

    def get_officer_profile(self, officer_id=None, role_filter=None):
        primary_file = os.path.join(DASHBOARD_DATA_DIR, "primary_learner.json")
        if not officer_id and not role_filter and os.path.exists(primary_file):
            with open(primary_file, "r") as f:
                return json.load(f)

        profiles_file = os.path.join(DATA_DIR, "official_profiles.json")
        if os.path.exists(profiles_file):
            with open(profiles_file, "r") as f:
                profiles = json.load(f)
                
            if officer_id:
                for p in profiles:
                    if p["officer_id"] == officer_id:
                        return p
                        
            if role_filter:
                for p in profiles:
                    if p.get("role_key") == role_filter:
                        return p
                        
            return profiles[0]
            
        return {}

    def search_officers(self, division=None, cadre=None, search=None, limit=50):
        profiles_file = os.path.join(DATA_DIR, "official_profiles.json")
        if not os.path.exists(profiles_file):
            return []

        with open(profiles_file, "r") as f:
            profiles = json.load(f)

        filtered = []
        for p in profiles:
            if division and division != "All" and p.get("division_code") != division:
                continue
            if cadre and cadre != "All" and cadre not in p.get("cadre", ""):
                continue
            if search:
                s_l = search.lower()
                if s_l not in p.get("name", "").lower() and s_l not in p.get("officer_id", "").lower() and s_l not in p.get("designation", "").lower():
                    continue
            filtered.append({
                "officer_id": p["officer_id"],
                "name": p["name"],
                "designation": p["designation"],
                "cadre": p["cadre"],
                "division_code": p["division_code"],
                "overall_competency_index": p["overall_competency_index"],
                "total_learning_hours": p["total_learning_hours"],
                "karma_points": p["karma_points"],
                "current_assignment": p["current_assignment"]
            })
            if len(filtered) >= limit:
                break
        return filtered

    def handle_course_enrolment(self, officer_id, course_id):
        primary_file = os.path.join(DASHBOARD_DATA_DIR, "primary_learner.json")
        if os.path.exists(primary_file):
            with open(primary_file, "r") as f:
                officer = json.load(f)
        else:
            return {"error": "Learner profile not found"}

        # Find course
        course = next((c for c in recommendation_engine.igot_courses if c["course_id"] == course_id), None)
        if not course:
            return {"error": "Course not found"}

        primary_comp = course["primary_competency"]
        cur_level = officer["current_competencies"].get(primary_comp, 2)
        tgt_level = officer["target_competencies"].get(primary_comp, 4)
        
        # Uplift level
        new_level = min(tgt_level, cur_level + 1)
        officer["current_competencies"][primary_comp] = new_level
        officer["total_learning_hours"] += course["duration_hours"]
        officer["karma_points"] += course["karma_points"]
        officer["completed_courses_count"] += 1
        
        # Ensure completed_courses array exists
        completed = officer.get("completed_courses", [])
        if course_id not in completed:
            completed.append(course_id)
        officer["completed_courses"] = completed
        
        # Recalculate gaps
        if primary_comp in officer["skill_gaps"]:
            officer["skill_gaps"][primary_comp]["current"] = new_level
            officer["skill_gaps"][primary_comp]["gap"] = max(0, tgt_level - new_level)
            officer["skill_gaps"][primary_comp]["severity"] = "High" if officer["skill_gaps"][primary_comp]["gap"] >= 2 else ("Medium" if officer["skill_gaps"][primary_comp]["gap"] == 1 else "None")

        # Uplift overall index
        officer["overall_competency_index"] = min(100.0, round(officer["overall_competency_index"] + 3.2, 1))

        # Save back
        with open(primary_file, "w") as f:
            json.dump(officer, f, indent=2)

        return {
            "success": True,
            "message": f"Successfully enrolled & completed certification for '{course['title']}'!",
            "new_competency_index": officer["overall_competency_index"],
            "karma_points_earned": course["karma_points"],
            "updated_officer": officer
        }

    def handle_nssta_nomination(self, officer_id, program_id):
        primary_file = os.path.join(DASHBOARD_DATA_DIR, "primary_learner.json")
        if os.path.exists(primary_file):
            with open(primary_file, "r") as f:
                officer = json.load(f)
        else:
            return {"error": "Learner profile not found"}

        # Track nominated programmes
        nominated = officer.get('nominated_programmes', [])
        if program_id not in nominated:
            nominated.append(program_id)
        officer['nominated_programmes'] = nominated
        
        # Save back
        with open(primary_file, "w") as f:
            json.dump(officer, f, indent=2)
            
        return {
            "success": True,
            "message": f"Successfully nominated for programme!",
            "updated_officer": officer
        }

    def apply_assessment_uplift(self, officer_id, eval_result):
        primary_file = os.path.join(DASHBOARD_DATA_DIR, "primary_learner.json")
        if not os.path.exists(primary_file):
            return

        with open(primary_file, "r") as f:
            officer = json.load(f)

        target_c = eval_result.get("target_competency", "STAT-01")
        uplift = eval_result.get("competency_level_uplift", 0.5)
        karma = eval_result.get("karma_points_awarded", 100)

        # Update Karma & Index
        officer["karma_points"] += karma
        officer["overall_competency_index"] = min(100.0, round(officer["overall_competency_index"] + uplift * 2.0, 1))
        
        # Uplift level
        if target_c in officer["current_competencies"]:
            cur = officer["current_competencies"][target_c]
            tgt = officer["target_competencies"].get(target_c, 4)
            officer["current_competencies"][target_c] = min(5, cur + (1 if uplift >= 0.5 else 0))
            if target_c in officer["skill_gaps"]:
                officer["skill_gaps"][target_c]["current"] = officer["current_competencies"][target_c]
                officer["skill_gaps"][target_c]["gap"] = max(0, tgt - officer["current_competencies"][target_c])

        with open(primary_file, "w") as f:
            json.dump(officer, f, indent=2)

    def handle_assistant_query(self, query: str, officer_id: str):
        q_lower = query.lower()
        
        if "gap" in q_lower or "skill" in q_lower or "deficiencies" in q_lower:
            return {
                "answer": "Based on your competency profile, your top priority skill gap is in **Modern Python for Microdata Processing (TECH-01)** (Gap: -2 Levels) and **AI & Machine Learning for Survey Nowcasting (TECH-07)**. Closing these gaps will elevate your Competency Readiness Index from **78.4% to 86.8%**.",
                "suggested_actions": ["Enrol in 'Python for Official Statistics (IGOT-TECH-201)'", "Take the Python Microdata Assessment", "Register for NSSTA-TPAC AI Workshop"],
                "recommended_course_id": "IGOT-TECH-201"
            }
        elif "national account" in q_lower or "sna" in q_lower or "gdp" in q_lower:
            return {
                "answer": "Under **SNA 2008** standards, Gross Value Added (GVA) at basic prices is defined as Gross Output minus Intermediate Consumption. GDP at market prices is derived by adding Net Product Taxes (Product Taxes - Product Subsidies). Supply-Use Tables (SUT) serve as the fundamental diagnostic tool for sector balancing.",
                "suggested_actions": ["Take the National Accounts Assessment", "Enrol in 'System of National Accounts & GDP (IGOT-STAT-102)'", "View SUT Module Syllabus"],
                "recommended_course_id": "IGOT-STAT-102"
            }
        elif "cpi" in q_lower or "price" in q_lower or "inflation" in q_lower:
            return {
                "answer": "The All-India CPI utilizes a **Modified Laspeyres price index** formula with base expenditure weights from the Household Consumption Expenditure Survey (HCES). At elementary aggregate levels, MoSPI mandates the **Jevons Index (geometric mean)** because it avoids the upward bias of Carli and satisfies the axiomatic time-reversal test.",
                "suggested_actions": ["Take CPI Compilation Quiz", "Explore 'Consumer Price Index & Inflation (IGOT-STAT-103)'"],
                "recommended_course_id": "IGOT-STAT-103"
            }
        elif "plfs" in q_lower or "labour" in q_lower or "employment" in q_lower:
            return {
                "answer": "The **Periodic Labour Force Survey (PLFS)** measures employment under two main criteria: **Usual Status (UPSS)** (major time >= 183 days + subsidiary work >= 30 days) and **Current Weekly Status (CWS)** (at least 1 hour of work on any 1 day during 7-day recall). In urban areas, a 25% rotational panel across 4 consecutive quarters is deployed.",
                "suggested_actions": ["Take PLFS Standards Quiz", "Enrol in 'PLFS Concepts & Data Analytics (IGOT-STAT-105)'"],
                "recommended_course_id": "IGOT-STAT-105"
            }
        elif "dpdpa" in q_lower or "privacy" in q_lower or "confidentiality" in q_lower:
            return {
                "answer": "Under the **Digital Personal Data Protection Act 2023 (DPDPA)**, MoSPI operates as a Significant Data Fiduciary. Before releasing public microdata, Statistical Disclosure Control (SDC) requires eliminating direct identifiers, enforcing **k-Anonymity (k >= 5)** on quasi-identifiers, and top-coding high-income outliers to prevent re-identification.",
                "suggested_actions": ["Take DPDPA 2023 Compliance Quiz", "Enrol in 'DPDPA 2023 & Microdata Privacy (IGOT-GOV-302)'"],
                "recommended_course_id": "IGOT-GOV-302"
            }
        else:
            return {
                "answer": f"I am your **KASHYAP AI Karmayogi Statistical Learning Assistant**. I can help you analyze your competency gaps, recommend iGOT Karmayogi micro-courses, navigate NSSTA TPAC in-service workshops, and explain statistical methodologies across SNA 2008, CPI, PLFS, Python for microdata, and DPDPA 2023.",
                "suggested_actions": ["Analyze my skill gaps", "Recommend iGOT courses for my role", "Generate an AI MCQ Quiz from a circular"],
                "recommended_course_id": "IGOT-STAT-101"
            }

def start_server():
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    socketserver.TCPServer.allow_reuse_address = True
    with socketserver.TCPServer(("", PORT), KashyapRequestHandler) as httpd:
        print(f"[OK] SKILL-STAT Web Server actively running at http://localhost:{PORT}/")
        try:
            httpd.serve_forever()
        except KeyboardInterrupt:
            print("\n[!] Shutting down SKILL-STAT Web Server.")

if __name__ == "__main__":
    start_server()
